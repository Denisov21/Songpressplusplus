###############################################################
# Name:             songimpress.py
# Purpose:     Generate PowerPoint presentation from songs
# Author:         Luca Allulli (webmaster@roma21.it)
# Modified by:  Denisov21
# Created:     2019-02-02
# Copyright: Luca Allulli (https://www.skeed.it/songpress)
#               Modifications copyright © 2026 Denisov21
# License:     GNU GPL v2
##############################################################

import sys

import wx
from pptx import Presentation

_ = wx.GetTranslation

# Layout usato per le slide dei canti, e gli idx dei placeholder delle due
# aree di testo al suo interno (riga corrente / riga successiva).
LAYOUT_INDEX = 3
BODY_IDX = (1, 2)


class SongPresentation:
    def __init__(self, template, out_filename):
        self.pres = Presentation(template)

        layouts = self.pres.slide_layouts
        if len(layouts) <= LAYOUT_INDEX:
            raise ValueError(_(
                "The template has no slide layout at index {index}: "
                "unable to generate the slides."
            ).format(index=LAYOUT_INDEX))
        self.layout = layouts[LAYOUT_INDEX]

        # Convalida per idx del placeholder, non per posizione nella lista: le
        # due aree di testo devono essere davvero presenti, altrimenti il testo
        # finirebbe silenziosamente nel placeholder sbagliato (es. piè di pagina
        # o data) oppure solleverebbe IndexError.
        available = {ph.placeholder_format.idx for ph in self.layout.placeholders}
        missing = [i for i in BODY_IDX if i not in available]
        if missing:
            raise ValueError(_(
                "Layout {index} of the template lacks the required text "
                "placeholders (idx {idx}): unable to generate the slides."
            ).format(index=LAYOUT_INDEX,
                     idx=', '.join(str(i) for i in missing)))

        self.prev = None
        self.out_filename = out_filename

    def _add_slide(self, cur, nxt):
        slide = self.pres.slides.add_slide(self.layout)
        ph = slide.placeholders          # accesso per idx, non per posizione nella lista
        ph[BODY_IDX[0]].text = cur
        ph[BODY_IDX[1]].text = nxt

    def _add_empty_slide(self):
        self.pres.slides.add_slide(self.layout)

    def add_line(self, line):
        if self.prev is not None:
            self._add_slide(self.prev, line)
        self.prev = line

    def end_song(self):
        if self.prev is not None:
            self._add_slide(self.prev, '')
        self._add_empty_slide()
        self.prev = None

    def close(self):
        if self.prev is not None:
            self._add_slide(self.prev, '')
        # Nessuna slide prodotta: l'input era vuoto o conteneva solo righe
        # bianche. Non salviamo un file a 0 slide (PowerPoint lo rifiuterebbe):
        # segnaliamo l'errore al chiamante.
        if len(self.pres.slides) == 0:
            raise ValueError(_(
                "No text to convert: the input contains no song lines."
            ))
        self.pres.save(self.out_filename)

    def __enter__(self):
        # Il blocco `with SongPresentation(...) as c` richiede __enter__:
        # senza, Python solleva TypeError ("does not support the context
        # manager protocol") prima ancora di eseguire il corpo.
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        if exc_type is not None:
            # Eccezione in corso: salva ciò che abbiamo ma segnalala al chiamante
            try:
                if self.prev is not None:
                    self._add_slide(self.prev, '')
                self.pres.save(self.out_filename)
            except Exception:
                pass
            return False   # rilancia l'eccezione originale
        self.close()
        return False


def to_presentation(lines, output_file, template_file):
    with SongPresentation(template_file, output_file) as c:
        for line in lines:
            line = line.replace('\n', '').strip()
            if line == '---':
                c.end_song()
            elif line != '':
                c.add_line(line)


if __name__ == '__main__':
    if len(sys.argv) != 4:
        print(_("Usage: songimpress.py SONG_FILE OUTPUT_PRESENTATION TEMPLATE_PRESENTATION"))
        sys.exit(1)

    with open(sys.argv[1], encoding='utf-8') as f:
        to_presentation(f, sys.argv[2], sys.argv[3])
